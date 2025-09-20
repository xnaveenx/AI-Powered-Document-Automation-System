import os
import docx
import pdfplumber
import openpyxl
from PIL import Image
import pytesseract
from pptx import Presentation
from backend.common.config import setting

# Configure Tesseract path
pytesseract.pytesseract.tesseract_cmd = setting.TESSERACT_PATH


def extract_pdf(file_path: str) -> dict:
    pages = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                pages.append(page_text)
            else:
                img = page.to_image(resolution=300).original
                ocr_text = pytesseract.image_to_string(img, lang=setting.TESSERACT_LANG)
                pages.append(ocr_text)
    return {"pages": pages, "full_text": "\n".join(pages)}


def extract_docx(file_path: str) -> dict:
    paras = [para.text for para in docx.Document(file_path).paragraphs if para.text.strip()]
    return {"pages": paras, "full_text": "\n".join(paras)}


def extract_xlsx(file_path: str) -> dict:
    wb = openpyxl.load_workbook(file_path)
    sheets_content = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        sheet_text = []
        for row in ws.iter_rows(values_only=True):
            row_text = " ".join([str(cell) for cell in row if cell])
            if row_text.strip():
                sheet_text.append(row_text)
        if sheet_text:
            sheets_content.append(f"[{sheet}] " + "\n".join(sheet_text))
    return {"pages": sheets_content, "full_text": "\n".join(sheets_content)}


def extract_txt(file_path: str) -> dict:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return {"pages": [text], "full_text": text}


def extract_image(file_path: str) -> dict:
    img = Image.open(file_path)
    ocr_text = pytesseract.image_to_string(img, lang=setting.TESSERACT_LANG)
    return {"pages": [ocr_text], "full_text": ocr_text}


def extract_pptx(file_path: str) -> dict:
    prs = Presentation(file_path)
    slides_content = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
            if hasattr(shape, "image"):
                img_bytes = shape.image.blob
                temp_img = "temp_slide.png"
                with open(temp_img, "wb") as f:
                    f.write(img_bytes)
                img = Image.open(temp_img)
                ocr_text = pytesseract.image_to_string(img, lang=setting.TESSERACT_LANG)
                if ocr_text.strip():
                    slide_text.append(ocr_text)
                os.remove(temp_img)
        slides_content.append("\n".join(slide_text))
    return {"pages": slides_content, "full_text": "\n".join(slides_content)}


def extract_any(file_path: str) -> dict:
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_pdf(file_path)
    elif ext == ".docx":
        return extract_docx(file_path)
    elif ext in [".xlsx", ".xls"]:
        return extract_xlsx(file_path)
    elif ext == ".txt":
        return extract_txt(file_path)
    elif ext in [".pptx", ".ppt"]:
        return extract_pptx(file_path)
    elif ext in [".png", ".jpg", ".jpeg", ".tiff"]:
        return extract_image(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")
